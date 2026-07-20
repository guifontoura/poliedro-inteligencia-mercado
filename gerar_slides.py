from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Inicializa a Apresentação Widescreen (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Paleta de Cores do Poliedro
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(2, 65, 161)  # Azul Escuro Forte
DARK_YELLOW = RGBColor(242, 193, 46)  # Amarelo Ouro
ORANGE = RGBColor(242, 115, 39)  # Laranja Accent
BLACK = RGBColor(33, 33, 33)  # Texto/Detalhes
LIGHT_GRAY = RGBColor(245, 247, 250)  # Fundo dos Cartões
GRAY = RGBColor(120, 120, 120)


def add_header(slide, title_text, category_text="NOVO PROGRAMA DE EXPANSÃO"):
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = 'Calibri'
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = ORANGE
    p_cat.space_after = Pt(4)

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.name = 'Calibri'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = BLUE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.4), Inches(11.83), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_YELLOW
    shape.line.fill.background()


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    return tf


def format_bullet(paragraph, bold_prefix, normal_text, font_size=16, color=BLACK):
    paragraph.space_after = Pt(12)
    paragraph.level = 0

    run1 = paragraph.add_run()
    run1.text = bold_prefix
    run1.font.bold = True
    run1.font.name = 'Calibri'
    run1.font.size = Pt(font_size)
    run1.font.color.rgb = color

    run2 = paragraph.add_run()
    run2.text = normal_text
    run2.font.bold = False
    run2.font.name = 'Calibri'
    run2.font.size = Pt(font_size)
    run2.font.color.rgb = color


def add_card(slide, left, top, width, height, bg_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


# --- SLIDE 1: CAPA ---
slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(slide_layout)
bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = BLUE
bg.line.fill.background()

accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), Inches(13.333), Inches(0.3))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = DARK_YELLOW
accent_bar.line.fill.background()

tf1 = add_textbox(slide1, Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
p1 = tf1.paragraphs[0]
p1.text = "MAPEAMENTO DE ESCOLAS PRIVADAS\nCOM ALTO POTENCIAL ECONÔMICO"
p1.font.name = 'Calibri'
p1.font.size = Pt(40)
p1.font.bold = True
p1.font.color.rgb = WHITE
p1.space_after = Pt(14)

p2 = tf1.add_paragraph()
p2.text = "Inteligência de Dados Aplicada à Expansão de Novos Negócios"
p2.font.name = 'Calibri'
p2.font.size = Pt(20)
p2.font.color.rgb = DARK_YELLOW
p2.space_after = Pt(40)

p3 = tf1.add_paragraph()
p3.text = "Guilherme Fontoura  |  Apresentação Executiva  |  Julho 2026"
p3.font.name = 'Calibri'
p3.font.size = Pt(14)
p3.font.color.rgb = WHITE

# --- SLIDE 2: O PROBLEMA ---
slide2 = prs.slides.add_slide(slide_layout)
add_header(slide2, "Onde estão nossos clientes ideais?")
add_card(slide2, Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8), LIGHT_GRAY)
tf2_left = add_textbox(slide2, Inches(1.1), Inches(2.1), Inches(4.9), Inches(4.2))
p_left_title = tf2_left.paragraphs[0]
p_left_title.text = "DESAFIOS DA PROSPECÇÃO TRADICIONAL"
p_left_title.font.size = Pt(16)
p_left_title.font.bold = True
p_left_title.font.color.rgb = ORANGE
p_left_title.space_after = Pt(20)
bullet1 = tf2_left.add_paragraph()
format_bullet(bullet1, "❌ Esforço Comercial Disperso: ",
              "Prospeccionar escolas privadas \"no escuro\" gera um custo de aquisição (CAC) alto e ineficiente.")
bullet2 = tf2_left.add_paragraph()
format_bullet(bullet2, "❌ Orçamentos Desiguais: ",
              "Nem toda escola particular tem margem ou perfil financeiro para adotar nossas soluções premium.")
bullet3 = tf2_left.add_paragraph()
format_bullet(bullet3, "❌ Ilusão do PIB Municipal: ",
              "Cidades ricas em PIB (indústrias, commodities) muitas vezes não têm famílias com alto poder de compra real.")

add_card(slide2, Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8), BLUE)
tf2_right = add_textbox(slide2, Inches(7.33), Inches(2.1), Inches(4.9), Inches(4.2))
p_right_title = tf2_right.paragraphs[0]
p_right_title.text = "A SOLUÇÃO INTELIGENTE DE MERCADO"
p_right_title.font.size = Pt(16)
p_right_title.font.bold = True
p_right_title.font.color.rgb = DARK_YELLOW
p_right_title.space_after = Pt(20)
bullet4 = tf2_right.add_paragraph()
format_bullet(bullet4, "✅ Cruzamento Estratégico: ",
              "Combinar o Censo Escolar (INEP) com a renda domiciliar do Censo (IBGE).", color=WHITE)
bullet5 = tf2_right.add_paragraph()
format_bullet(bullet5, "✅ Foco no Bolso dos Pais: ",
              "Identificar com precisão cirúrgica as famílias que de fato pagam as mensalidades.", color=WHITE)
bullet6 = tf2_right.add_paragraph()
format_bullet(bullet6, "✅ Segmentação Adequada: ",
              "Direcionar o produto certo (Poliedro vs. Polígono) para o perfil econômico da região.", color=WHITE)

# --- SLIDE 3: METODOLOGIA ---
slide3 = prs.slides.add_slide(slide_layout)
add_header(slide3, "Método de Ranking: Três Variáveis Inteligentes")
box_y = Inches(2.0)
box_height = Inches(4.5)
box_width = Inches(3.64)
gap = Inches(0.45)
add_card(slide3, Inches(0.75), box_y, box_width, box_height, LIGHT_GRAY, border_color=BLUE)
tf3_1 = add_textbox(slide3, Inches(1.0), Inches(2.3), Inches(3.14), Inches(4.0))
p_b1_title = tf3_1.paragraphs[0]
p_b1_title.text = "1️⃣ RENDA DOMICILIAR REAL\n(Fator Principal)"
p_b1_title.font.size = Pt(16)
p_b1_title.font.bold = True
p_b1_title.font.color.rgb = BLUE
p_b1_title.space_after = Pt(14)
p_b1_desc1 = tf3_1.add_paragraph()
format_bullet(p_b1_desc1, "Métrica: ",
              "Percentual de famílias com rendimento acima de 5 salários mínimos (Censo 2022).", font_size=14)
p_b1_desc2 = tf3_1.add_paragraph()
format_bullet(p_b1_desc2, "Por quê: ",
              "Indica o poder de compra real e define o fit para o Poliedro Premium ou para a flexibilidade do Polígono.",
              font_size=14)

add_card(slide3, Inches(0.75) + box_width + gap, box_y, box_width, box_height, LIGHT_GRAY, border_color=DARK_YELLOW)
tf3_2 = add_textbox(slide3, Inches(1.0) + box_width + gap, Inches(2.3), Inches(3.14), Inches(4.0))
p_b2_title = tf3_2.paragraphs[0]
p_b2_title.text = "2️⃣ POPULAÇÃO INFANTIL\n(Demanda Ativa)"
p_b2_title.font.size = Pt(16)
p_b2_title.font.bold = True
p_b2_title.font.color.rgb = BLACK
p_b2_title.space_after = Pt(14)
p_b2_desc1 = tf3_2.add_paragraph()
format_bullet(p_b2_desc1, "Métrica: ", "População de 0 a 17 anos residente no município (Censo 2022).", font_size=14)
p_b2_desc2 = tf3_2.add_paragraph()
format_bullet(p_b2_desc2, "Filtro Técnico: ",
              "Mínimo de 1.000 crianças para elegibilidade. Cidades ricas envelhecidas são mercados sem futuro.",
              font_size=14)

add_card(slide3, Inches(0.75) + (box_width + gap) * 2, box_y, box_width, box_height, LIGHT_GRAY, border_color=ORANGE)
tf3_3 = add_textbox(slide3, Inches(1.0) + (box_width + gap) * 2, Inches(2.3), Inches(3.14), Inches(4.0))
p_b3_title = tf3_3.paragraphs[0]
p_b3_title.text = "3️⃣ PORTE E INFRAESTRUTURA\n(Incentivo ao Orçamento)"
p_b3_title.font.size = Pt(16)
p_b3_title.font.bold = True
p_b3_title.font.color.rgb = ORANGE
p_b3_title.space_after = Pt(14)
p_b3_desc1 = tf3_3.add_paragraph()
format_bullet(p_b3_desc1, "Métrica: ", "Número de salas de aula ativas (INEP Censo Escolar 2023).", font_size=14)
p_b3_desc2 = tf3_3.add_paragraph()
format_bullet(p_b3_desc2, "Por quê: ",
              "Escolas maiores possuem melhor infraestrutura e maior verba disponível para fechamento de contratos complexos.",
              font_size=14)

# --- SLIDE 4: ESCOPO DO PROJETO ---
slide4 = prs.slides.add_slide(slide_layout)
add_header(slide4, "Escopo e Resultados da Análise")
tf4_left = add_textbox(slide4, Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.5))
p_sc_title = tf4_left.paragraphs[0]
p_sc_title.text = "DADOS TRATADOS E ANALISADOS"
p_sc_title.font.size = Pt(18)
p_sc_title.font.bold = True
p_sc_title.font.color.rgb = BLUE
p_sc_title.space_after = Pt(20)


def add_stat(slide, left, top, value, label, text_color=BLUE):
    box = slide.shapes.add_textbox(left, top, Inches(5.5), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p_val = tf.paragraphs[0]
    p_val.text = value
    p_val.font.name = 'Calibri'
    p_val.font.size = Pt(36)
    p_val.font.bold = True
    p_val.font.color.rgb = text_color
    p_val.space_after = Pt(2)
    p_lbl = tf.add_paragraph()
    p_lbl.text = label
    p_lbl.font.name = 'Calibri'
    p_lbl.font.size = Pt(12)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = BLACK


add_stat(slide4, Inches(0.75), Inches(2.5), "32.959", "Escolas privadas mapeadas nacionalmente (INEP)")
add_stat(slide4, Inches(0.75), Inches(4.0), "3.100+", "Municípios mapeados e analisados")
add_stat(slide4, Inches(0.75), Inches(5.5), "755", "Municípios elegíveis com alta densidade infantil ativa")
add_stat(slide4, Inches(6.98), Inches(2.5), "0.42 / 1.00", "Score médio geral de potencial das escolas")
add_stat(slide4, Inches(6.98), Inches(4.0), "0.99", "Score das escolas Elite no Top 1% nacional")
add_stat(slide4, Inches(6.98), Inches(5.5), "11.2", "Média nacional de salas de aula por escola")

# --- SLIDE 5: TOP 10 OPORTUNIDADES ---
slide5 = prs.slides.add_slide(slide_layout)
add_header(slide5, "Os Melhores Mercados Nacionais para Prospecção")
rows, cols = 6, 5
left, top, width, height = Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.2)
table_shape = slide5.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table
table.columns[0].width = Inches(0.5)
table.columns[1].width = Inches(4.5)
table.columns[2].width = Inches(3.0)
table.columns[3].width = Inches(1.83)
table.columns[4].width = Inches(2.0)
headers = ["#", "Escola de Alto Potencial", "Município (UF)", "Score Geral", "Salas Ativas"]
for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if i != 1 else PP_ALIGN.LEFT
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.font.size = Pt(13)
data = [
    ["1", "Centro Educacional Primeiro Mundo", "Canaã dos Carajás (PA)", "0.993", "31 salas"],
    ["2", "Internacional de Alphaville", "Barueri (SP)", "0.991", "24 salas"],
    ["3", "Sant Anna Colegio", "Campinas (SP)", "0.990", "24 salas"],
    ["4", "Colegio Bom Jesus", "Itajaí (SC)", "0.986", "62 salas"],
    ["5", "Nossa Senhora de Misericórdia", "Blumenau (SC)", "0.984", "40 salas"]
]
for row_idx, row_data in enumerate(data):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = value
        cell.fill.solid()
        if (row_idx + 1) % 2 == 0:
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.fore_color.rgb = WHITE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if col_idx != 1 else PP_ALIGN.LEFT
            p.font.color.rgb = BLACK
            p.font.name = 'Calibri'
            p.font.size = Pt(12)
note_box = slide5.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(11.83), Inches(0.5))
p_note = note_box.text_frame.paragraphs[0]
p_note.text = "💡 Os líderes do ranking combinam renda de topo de pirâmide, excelente tamanho de infraestrutura e mercado infantil saudável."
p_note.font.italic = True
p_note.font.size = Pt(12)
p_note.font.color.rgb = ORANGE

# --- SLIDE 6: SEGMENTAÇÃO ---
slide6 = prs.slides.add_slide(slide_layout)
add_header(slide6, "Três Segmentos Estratégicos de Oportunidade")
add_card(slide6, Inches(0.75), box_y, box_width, box_height, LIGHT_GRAY, border_color=BLUE)
tf6_1 = add_textbox(slide6, Inches(1.0), Inches(2.3), Inches(3.14), Inches(4.0))
p6_1 = tf6_1.paragraphs[0]
p6_1.text = "🥇 GOLDEN LEADS\n(Score 0.90 a 1.00)"
p6_1.font.size = Pt(16)
p6_1.font.bold = True
p6_1.font.color.rgb = BLUE
p6_1.space_after = Pt(14)
p6_1_b1 = tf6_1.add_paragraph()
format_bullet(p6_1_b1, "Foco: ", "Poliedro Tradicional (Alta Renda).", font_size=13)
p6_1_b2 = tf6_1.add_paragraph()
format_bullet(p6_1_b2, "Volume: ", "150 escolas de alta infraestrutura.", font_size=13)
p6_1_b3 = tf6_1.add_paragraph()
format_bullet(p6_1_b3, "Ação Comercial: ", "Proposta consultiva premium direta com mantenedores.", font_size=13)

add_card(slide6, Inches(0.75) + box_width + gap, box_y, box_width, box_height, LIGHT_GRAY, border_color=DARK_YELLOW)
tf6_2 = add_textbox(slide6, Inches(1.0) + box_width + gap, Inches(2.3), Inches(3.14), Inches(4.0))
p6_2 = tf6_2.paragraphs[0]
p6_2.text = "🥈 PROSPECTS QUALIF.\n(Score 0.70 a 0.89)"
p6_2.font.size = Pt(16)
p6_2.font.bold = True
p6_2.font.color.rgb = BLACK
p6_2.space_after = Pt(14)
p6_2_b1 = tf6_2.add_paragraph()
format_bullet(p6_2_b1, "Foco: ", "Polígono (Classe Média em crescimento).", font_size=13)
p6_2_b2 = tf6_2.add_paragraph()
format_bullet(p6_2_b2, "Volume: ", "1.200 escolas com boa estrutura local.", font_size=13)
p6_2_b3 = tf6_2.add_paragraph()
format_bullet(p6_2_b3, "Ação Comercial: ", "Inbound estruturado focado em custo x benefício e flexibilidade.",
              font_size=13)

add_card(slide6, Inches(0.75) + (box_width + gap) * 2, box_y, box_width, box_height, LIGHT_GRAY, border_color=ORANGE)
tf6_3 = add_textbox(slide6, Inches(1.0) + (box_width + gap) * 2, Inches(2.3), Inches(3.14), Inches(4.0))
p6_3 = tf6_3.paragraphs[0]
p6_3.text = "🥉 PIPELINE FUTURO\n(Score 0.40 a 0.69)"
p6_3.font.size = Pt(16)
p6_3.font.bold = True
p6_3.font.color.rgb = ORANGE
p6_3.space_after = Pt(14)
p6_3_b1 = tf6_3.add_paragraph()
format_bullet(p6_3_b1, "Foco: ", "Nutrição de marca a longo prazo.", font_size=13)
p6_3_b2 = tf6_3.add_paragraph()
format_bullet(p6_3_b2, "Volume: ", "Mais de 31.000 escolas emergentes.", font_size=13)
p6_3_b3 = tf6_3.add_paragraph()
format_bullet(p6_3_b3, "Ação Comercial: ", "Automação de e-mails, webinars pedagógicos e e-books.", font_size=13)

# --- SLIDE 7: COMPARATIVO ---
slide7 = prs.slides.add_slide(slide_layout)
add_header(slide7, "O Salto de Eficiência Comercial")
rows_7, cols_7 = 5, 3
t7_left, t7_top, t7_width, t7_height = Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.2)
table_shape7 = slide7.shapes.add_table(rows_7, cols_7, t7_left, t7_top, t7_width, t7_height)
table7 = table_shape7.table
table7.columns[0].width = Inches(3.0)
table7.columns[1].width = Inches(4.415)
table7.columns[2].width = Inches(4.415)
headers7 = ["Critério", "Antes (Prospecção Tradicional)", "Depois (Mapeamento por Dados)"]
for i, header in enumerate(headers7):
    cell = table7.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE if i > 0 else GRAY
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.font.size = Pt(13)
data7 = [
    ["Encontrar clientes", "Listas de escolas aleatórias, ligando \"frio\"",
     "Identificação cirúrgica por Score de Potencial"],
    ["Previsão de orçamento", "\"Achismo\" ou tentativa e erro em visitas",
     "Validação direta da renda familiar real via Censo"],
    ["Tempo gasto", "Semanas ou meses de qualificação manual", "Pronto em segundos na planilha filtrada"],
    ["Conversão esperada", "Baixa (< 2%) pela dispersão do esforço", "Elevada, focando apenas no topo de pirâmide"]
]
for row_idx, row_data in enumerate(data7):
    for col_idx, value in enumerate(row_data):
        cell = table7.cell(row_idx + 1, col_idx)
        cell.text = value
        cell.fill.solid()
        if (row_idx + 1) % 2 == 0:
            cell.fill.fore_color.rgb = LIGHT_GRAY
        else:
            cell.fill.fore_color.rgb = WHITE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            p.font.color.rgb = BLACK
            p.font.name = 'Calibri'
            p.font.size = Pt(12)

# --- SLIDE 8: PRÓXIMOS PASSOS ---
slide8 = prs.slides.add_slide(slide_layout)
add_header(slide8, "Plano de Implementação Prática")
add_card(slide8, Inches(0.75), box_y, box_width, box_height, LIGHT_GRAY)
tf8_1 = add_textbox(slide8, Inches(1.0), Inches(2.3), Inches(3.14), Inches(4.0))
p8_1 = tf8_1.paragraphs[0]
p8_1.text = "🎯 FASE 1: VALIDAÇÃO\n(Semanas 1 e 2)"
p8_1.font.size = Pt(16)
p8_1.font.bold = True
p8_1.font.color.rgb = BLUE
p8_1.space_after = Pt(14)
p8_1_b1 = tf8_1.add_paragraph()
format_bullet(p8_1_b1, "Atividade: ", "Selecionar 10 Golden Leads do top 50.", font_size=13)
p8_1_b2 = tf8_1.add_paragraph()
format_bullet(p8_1_b2, "Foco: ", "Realizar abordagem de vendas experimental e qualificar o retorno.", font_size=13)
p8_1_b3 = tf8_1.add_paragraph()
format_bullet(p8_1_b3, "Objetivo: ", "Comprovar a aderência do score na prática.", font_size=13)

add_card(slide8, Inches(0.75) + box_width + gap, box_y, box_width, box_height, LIGHT_GRAY)
tf8_2 = add_textbox(slide8, Inches(1.0) + box_width + gap, Inches(2.3), Inches(3.14), Inches(4.0))
p8_2 = tf8_2.paragraphs[0]
p8_2.text = "🚀 FASE 2: PILOTO\n(Mês 1)"
p8_2.font.size = Pt(16)
p8_2.font.bold = True
p8_2.font.color.rgb = DARK_YELLOW
p8_2.space_after = Pt(14)
p8_2_b1 = tf8_2.add_paragraph()
format_bullet(p8_2_b1, "Atividade: ", "Expandir para os 150 Golden Leads.", font_size=13)
p8_2_b2 = tf8_2.add_paragraph()
format_bullet(p8_2_b2, "Suporte: ", "Munir consultores com insights demográficos locais de cada escola.", font_size=13)
p8_2_b3 = tf8_2.add_paragraph()
format_bullet(p8_2_b3, "Objetivo: ", "Gerar os primeiros fechamentos com alta conversão.", font_size=13)

add_card(slide8, Inches(0.75) + (box_width + gap) * 2, box_y, box_width, box_height, LIGHT_GRAY)
tf8_3 = add_textbox(slide8, Inches(1.0) + (box_width + gap) * 2, Inches(2.3), Inches(3.14), Inches(4.0))
p8_3 = tf8_3.paragraphs[0]
p8_3.text = "📈 FASE 3: ESCALA\n(Mês 2+)"
p8_3.font.size = Pt(16)
p8_3.font.bold = True
p8_3.font.color.rgb = ORANGE
p8_3.space_after = Pt(14)
p8_3_b1 = tf8_3.add_paragraph()
format_bullet(p8_3_b1, "Atividade: ", "Automação e importação das notas no CRM.", font_size=13)
p8_3_b2 = tf8_3.add_paragraph()
format_bullet(p8_3_b2, "Segmento 2: ", "Iniciar prospecção de massa para as 1.200 escolas do portfólio Polígono.",
              font_size=13)
p8_3_b3 = tf8_3.add_paragraph()
format_bullet(p8_3_b3, "Objetivo: ", "Industrializar a inteligência de expansão.", font_size=13)

# --- SLIDE 9: EVOLUÇÃO ---
slide9 = prs.slides.add_slide(slide_layout)
add_header(slide9, "Roadmap Técnico de Evolução", category_text="VISÃO DE PRODUTO")
tf9 = add_textbox(slide9, Inches(0.75), Inches(2.0), Inches(11.83), Inches(4.5))
bullet9_1 = tf9.paragraphs[0]
format_bullet(bullet9_1, "⚡ Versão 1.0 (Atual): ",
              "Score de potencial estático gerado e pronto em planilhas unificando Censo + INEP.")
bullet9_2 = tf9.add_paragraph()
format_bullet(bullet9_2, "💻 Versão 2.0: ",
              "Desenvolvimento de um painel interativo (Streamlit) para filtros personalizados por UF/Município.")
bullet9_3 = tf9.add_paragraph()
format_bullet(bullet9_3, "🔌 Versão 3.0: ",
              "Integração via API com o CRM (HubSpot/Salesforce) para auto-scoring de novos cadastros de leads.")
bullet9_4 = tf9.add_paragraph()
format_bullet(bullet9_4, "🗺️ Versão 4.0: ",
              "Georreferenciamento de concorrência e visualização espacial das escolas em mapas dinâmicos.")
bullet9_5 = tf9.add_paragraph()
format_bullet(bullet9_5, "🤖 Versão 5.0: ", "Modelo preditivo de propensão de fechamento integrado ao funil de vendas.")

# --- SLIDE 10: CONCLUSÃO ---
slide10 = prs.slides.add_slide(slide_layout)
bg10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg10.fill.solid()
bg10.fill.fore_color.rgb = BLUE
bg10.line.fill.background()
accent_bar10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), Inches(13.333), Inches(0.3))
accent_bar10.fill.solid()
accent_bar10.fill.fore_color.rgb = DARK_YELLOW
accent_bar10.line.fill.background()
tf10 = add_textbox(slide10, Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.0))
p10_title = tf10.paragraphs[0]
p10_title.text = "PRONTOS PARA CRESCER COM INTELIGÊNCIA"
p10_title.font.size = Pt(36)
p10_title.font.bold = True
p10_title.font.color.rgb = WHITE
p10_title.space_after = Pt(24)
p10_b1 = tf10.add_paragraph()
format_bullet(p10_b1, "🎯 Eficiência: ",
              "Redução drástica das horas de qualificação comercial para focar em quem pode comprar.", font_size=18,
              color=WHITE)
p10_b2 = tf10.add_paragraph()
format_bullet(p10_b2, "🎯 Casamento de Portfólio: ",
              "Poliedro tradicional para o topo financeiro; Polígono para o volume da classe média.", font_size=18,
              color=WHITE)
p10_b3 = tf10.add_paragraph()
format_bullet(p10_b3, "🎯 Praticidade: ", "Framework pronto para testes em até 3 semanas se fizer sentido para o time.",
              font_size=18, color=WHITE)
p10_end = tf10.add_paragraph()
p10_end.text = "\n✋ Dúvidas ou sugestões?\nObrigado!"
p10_end.font.size = Pt(20)
p10_end.font.bold = True
p10_end.font.color.rgb = DARK_YELLOW

# Salva
prs.save("Mapeamento_Escolas_Poliedro.pptx")
print("Arquivo salvo localmente com sucesso!")